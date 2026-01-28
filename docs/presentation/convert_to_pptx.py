#!/usr/bin/env python3
"""
Convert Markdown presentation to PowerPoint format
Improved version with better slide structure handling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown(file_path):
    """Parse markdown file and extract slides content"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by --- separator
    sections = content.split('---')

    slides = []
    title_slide_done = False

    for section in sections:
        if not section.strip():
            continue

        lines = section.strip().split('\n')
        slide_data = {'title': '', 'subtitles': [], 'content': [], 'table': None, 'code': []}
        in_code_block = False
        code_lines = []
        table_lines = []

        for line in lines:
            line = line.rstrip('\n')

            # Skip empty lines
            if not line.strip():
                if in_code_block:
                    code_lines.append('')
                continue

            # Code block
            if line.strip().startswith('```'):
                if in_code_block:
                    # End code block
                    slide_data['code'] = code_lines
                    code_lines = []
                    in_code_block = False
                else:
                    # Start code block
                    in_code_block = True
                continue

            if in_code_block:
                code_lines.append(line)
                continue

            # Table
            if line.strip().startswith('|'):
                table_lines.append(line)
                continue

            # Main title (# Title) - only for title slide
            if line.startswith('# ') and not line.startswith('##'):
                if not title_slide_done:
                    slide_data['title'] = line[2:].strip()
                continue

            # Section title (## Title) or Slide title
            if line.startswith('## '):
                slide_data['title'] = line[3:].strip()
                continue

            # Subtitle (### Title)
            if line.startswith('### '):
                slide_data['subtitles'].append(line[4:].strip())
                continue

            # Regular content
            slide_data['content'].append(line.strip())

        # Process table if exists
        if table_lines:
            table_data = []
            for table_line in table_lines:
                if '---' not in table_line:  # Skip separator line
                    cells = [cell.strip() for cell in table_line.split('|')[1:-1]]
                    if cells:
                        table_data.append(cells)
            if table_data:
                slide_data['table'] = table_data

        # Only add slide if it has content or title
        if slide_data['title'] or slide_data['content'] or slide_data['table'] or slide_data['code']:
            if not title_slide_done:
                title_slide_done = True
            slides.append(slide_data)

    return slides

def add_title_slide(prs, title, content_lines):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 136, 229)  # #1E88E5

    # Subtitle info
    if content_lines:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True

        for i, line in enumerate(content_lines):
            if i > 0:
                subtitle_frame.add_paragraph()
            p = subtitle_frame.paragraphs[i]
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(66, 165, 245)  # #42A5F5
            p.space_after = Pt(6)

def add_section_slide(prs, title):
    """Add section separator slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 136, 229)

def add_content_slide(prs, slide_data):
    """Add content slide with title, content, table, and code"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = slide_data['title']
    subtitles = slide_data['subtitles']
    content = slide_data['content']
    table = slide_data['table']
    code = slide_data['code']

    # Calculate vertical positioning
    top_pos = Inches(0.5)

    # Title
    if title:
        title_box = slide.shapes.add_textbox(Inches(0.5), top_pos, Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title

        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(30)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 136, 229)

        top_pos += Inches(0.9)

    # Subtitles
    if subtitles:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), top_pos, Inches(9), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True

        for i, subtitle in enumerate(subtitles):
            if i > 0:
                subtitle_frame.add_paragraph()

            p = subtitle_frame.paragraphs[i]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(66, 165, 245)
            p.space_after = Pt(4)

        top_pos += Inches(0.7)

    # Content
    has_table_or_code = bool(table or code)
    content_height = Inches(2.5) if has_table_or_code else Inches(5.5)

    if content and not has_table_or_code:
        content_box = slide.shapes.add_textbox(Inches(0.5), top_pos, Inches(9), content_height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, text in enumerate(content):
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.text = format_text(text)
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(33, 33, 33)
            p.space_after = Pt(8)
            p.line_spacing = 1.4

        top_pos += content_height
    elif content:
        # Short content when there's table/code
        content_box = slide.shapes.add_textbox(Inches(0.5), top_pos, Inches(9), Inches(1.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, text in enumerate(content[:5]):  # Limit to 5 lines
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.text = format_text(text)
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(33, 33, 33)
            p.space_after = Pt(6)

        top_pos += Inches(1.6)

    # Table
    if table and len(table) > 0:
        rows = min(len(table), 10)  # Limit to 10 rows
        cols = len(table[0])

        table_left = Inches(0.5)
        table_top = top_pos
        table_width = Inches(9)
        table_height = Inches(3)

        shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
        tbl = shape.table

        # Adjust column widths
        for col in range(cols):
            tbl.columns[col].width = Inches(table_width.inches / cols)

        for r in range(rows):
            for c in range(cols):
                cell = tbl.cell(r, c)
                cell.text = str(table[r][c]) if r < len(table) and c < len(table[r]) else ''

                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11 if rows > 8 else 12)
                para.alignment = PP_ALIGN.CENTER

                # Header row
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(30, 136, 229)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                # Alternate rows
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 248, 255)

        top_pos += Inches(3.2)

    # Code block
    if code:
        code_top = top_pos if not table else Inches(4.5)
        code_height = Inches(2.5)

        # Background
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.4), code_top - Inches(0.05),
            Inches(9.2), code_height + Inches(0.1)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(245, 245, 245)
        background.line.color.rgb = RGBColor(200, 200, 200)

        code_box = slide.shapes.add_textbox(Inches(0.5), code_top, Inches(9), code_height)
        code_frame = code_box.text_frame
        code_frame.word_wrap = True

        # Join code lines
        code_text = '\n'.join(code[:25])  # Limit to 25 lines
        code_frame.text = code_text

        for paragraph in code_frame.paragraphs:
            paragraph.font.name = 'Consolas'
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(50, 50, 50)
            paragraph.space_after = Pt(0)

def format_text(text):
    """Format markdown text with basic formatting"""
    # Remove markdown bold markers for now (keep text)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'__(.*?)__', r'\1', text)

    # Handle checkboxes
    text = re.sub(r'- \[x\]', '✅', text)
    text = re.sub(r'- \[ \]', '⬜', text)

    return text

def create_presentation(md_file, output_file):
    """Create PowerPoint presentation from markdown"""
    print("📖 Parsing markdown file...")
    slides_data = parse_markdown(md_file)

    print(f"📊 Found {len(slides_data)} slides")
    print("🎨 Creating PowerPoint presentation...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Process slides
    for i, slide_data in enumerate(slides_data, 1):
        title = slide_data['title']

        # Check if it's a section slide (has # in title and little content)
        is_section = (title.startswith('第') and '部分' in title and
                      not slide_data['content'] and
                      not slide_data['table'] and
                      not slide_data['code'])

        # First slide is title slide
        if i == 1:
            add_title_slide(prs, title, slide_data['content'])
        elif is_section:
            add_section_slide(prs, title)
        else:
            add_content_slide(prs, slide_data)

        if i % 10 == 0:
            print(f"  Processed {i}/{len(slides_data)} slides...")

    # Save
    prs.save(output_file)
    print(f"\n✅ Presentation saved successfully!")
    print(f"   Location: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"\n🎓 Ready for your defense presentation!")

if __name__ == '__main__':
    md_file = '/Users/wangbo/gocache/docs/presentation/答辩PPT.md'
    output_file = '/Users/wangbo/gocache/docs/presentation/GoCache答辩.pptx'

    create_presentation(md_file, output_file)
